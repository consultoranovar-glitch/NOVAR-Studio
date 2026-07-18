#!/usr/bin/env python3
"""
Motor de presentaciones NOVAR · variantes desde una BASE ÚNICA.

Principio (regla de oro del estudio): cuando existe un PPTX canónico validado,
las variantes NO se reconstruyen — se CLONAN reemplazando solo el texto indicado,
preservando el formato del primer run de cada forma.

Uso: python3 base_unica_swap.py base.pptx mapa.json salida.pptx

mapa.json: lista de reemplazos [{"slide": 1, "shape": 5, "texto": "..."}]
(slide 1-based, shape = índice en slide.shapes; "\n" crea párrafos nuevos).
Para descubrir índices: python3 base_unica_swap.py base.pptx --inspeccionar
"""
import json, sys
from pptx import Presentation

def copy_font(src, dst):
    f, g = src.font, dst.font
    try: g.size = f.size
    except Exception: pass
    g.bold = f.bold; g.italic = f.italic
    if f.name: g.name = f.name
    try:
        if f.color and f.color.type is not None: g.color.rgb = f.color.rgb
    except Exception: pass

def set_text(shape, newtext):
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    r0 = p0.runs[0] if p0.runs else p0.add_run()
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    lines = newtext.split("\n")
    r0.text = lines[0]
    for ln in lines[1:]:
        np = tf.add_paragraph(); np.alignment = p0.alignment
        nr = np.add_run(); nr.text = ln; copy_font(r0, nr)

def inspeccionar(base):
    prs = Presentation(base)
    for si, s in enumerate(prs.slides, 1):
        for hi, sh in enumerate(s.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip():
                print(f"slide {si} · shape {hi}: {sh.text_frame.text.strip()[:80]!r}")

def main():
    if len(sys.argv) >= 3 and sys.argv[2] == "--inspeccionar":
        inspeccionar(sys.argv[1]); return
    if len(sys.argv) != 4:
        sys.exit(__doc__)
    base, mapa, out = sys.argv[1:4]
    prs = Presentation(base)
    slides = list(prs.slides)
    aplicados = 0
    for item in json.load(open(mapa, encoding="utf-8")):
        shapes = list(slides[item["slide"] - 1].shapes)
        sh = shapes[item["shape"]]
        if sh.has_text_frame:
            set_text(sh, item["texto"]); aplicados += 1
        else:
            print("AVISO: la forma no tiene texto:", item["slide"], item["shape"])
    prs.save(out)
    print(f"{aplicados} reemplazos · {out}")

if __name__ == "__main__":
    main()
